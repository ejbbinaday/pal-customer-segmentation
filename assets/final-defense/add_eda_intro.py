"""Insert a Problem / Solution / Limitations framing slide at the head of the EDA section.

Matches the deck's existing system: Arial throughout, #0D2E5E navy ink, #2166C0 eyebrow, and the
accent ladder #0D2E5E -> #1B74E8 -> #35A0DB for the three column rules.

The lower half carries a shape-built diagram of the grain change — a real four-coupon round trip
via Manila collapsing into the single booking that represents one purchase decision. Everything is
native PowerPoint geometry, so it stays editable and needs no linked image.
"""
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SRC, DST, LOGO = sys.argv[1], sys.argv[2], sys.argv[3]

NAVY    = RGBColor(0x0D, 0x2E, 0x5E)   # ink + headings
EYEBROW = RGBColor(0x21, 0x66, 0xC0)   # section label
BRIGHT  = RGBColor(0x1B, 0x74, 0xE8)   # accent rail 1
LIGHT   = RGBColor(0x35, 0xA0, 0xDB)   # accent rail 2
GREY    = RGBColor(0x6B, 0x7A, 0x8C)   # captions / closer
CHIP_BG = RGBColor(0xEC, 0xF2, 0xFA)   # coupon chip fill
CHIP_LN = RGBColor(0xC5, 0xD8, 0xEE)   # coupon chip hairline
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(SRC)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank


def textbox(left, top, width, height, runs, size, color=NAVY, bold=False,
            spacing=None, align=PP_ALIGN.LEFT, space_pt=0.0):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    for text, is_bold in runs:
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold or is_bold
        r.font.color.rgb = color
        if space_pt:
            r.font._rPr.set("spc", str(int(space_pt * 100)))
    return tb


def bar(left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def chip(left, top, width, height, text, fill, ink, size, bold=False, radius=0.14):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if fill is CHIP_BG:
        sh.line.color.rgb = CHIP_LN
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = ink
    return sh


# ── slide chrome, cloned from the other EDA slides ────────────────────────────
textbox(1.0, 0.42, 9.5, 0.3,
        [("Market Segmentation for Philippine Airlines Through Machine Learning", True)], 10.5)
bar(1.0, 0.8, 11.33, 0.02, NAVY)
slide.shapes.add_picture(LOGO, Inches(11.72), Inches(0.3), Inches(0.57), Inches(0.34))
textbox(1.0, 1.06, 11.33, 0.3, [("EDA  ·  JOSH", True)], 11.0, EYEBROW)
textbox(0.98, 1.36, 11.33, 1.0,
        [("A coupon is one journey; a booking is one decision", True)], 27.0)

# ── the three columns ─────────────────────────────────────────────────────────
COL_W, GAP, TOP = 3.55, 0.34, 2.30
cols = [
    ("Problem", NAVY, [
        ("The given data set ", False),
        ("was on a coupon level (a single journey)", True),
        (" which is not representative of the booker's decision", False)]),
    ("Solution", BRIGHT, [
        ("Rolling-up the data into the ", False),
        ("ticketed booking level", True),
        (", using the unique ID field.", False)]),
    ("Limitations", LIGHT, [
        ("We were not allowed to use PII.", False)]),
]
for i, (head, colour, body) in enumerate(cols):
    left = 1.0 + i * (COL_W + GAP)
    bar(left, TOP, COL_W, 0.055, colour)
    textbox(left, TOP + 0.24, COL_W, 0.34, [(head, True)], 16.0)
    textbox(left, TOP + 0.74, COL_W, 1.4, body, 13.0, spacing=1.25)

# ── the diagram: four coupons of one trip becoming one booking ────────────────
DTOP = 4.22
textbox(1.0, DTOP, 5.1, 0.24,
        [("AS GIVEN  ·  COUPON LEVEL", True)], 10.0, GREY, space_pt=0.8)
textbox(7.35, DTOP, 4.98, 0.24,
        [("AS MODELLED  ·  BOOKING LEVEL", True)], 10.0, EYEBROW, space_pt=0.8)

# one traveller's round trip via Manila = four separate rows in the raw extract
LEGS = ["CEB → MNL", "MNL → DXB", "DXB → MNL", "MNL → CEB"]
CW, CG = 1.17, 0.14
for i, leg in enumerate(LEGS):
    chip(1.0 + i * (CW + CG), DTOP + 0.34, CW, 0.62, leg, CHIP_BG, NAVY, 10.5)
textbox(1.0, DTOP + 1.06, 5.1, 0.5,
        [("Four rows. Counting these as four customers over-counts the "
          "cheapest journeys and hides the decision.", False)], 10.5, GREY, spacing=1.2)

arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(6.28), Inches(DTOP + 0.49), Inches(0.8), Inches(0.32))
arrow.fill.solid()
arrow.fill.fore_color.rgb = BRIGHT
arrow.line.fill.background()
arrow.shadow.inherit = False
textbox(5.95, DTOP + 0.90, 1.45, 0.24,
        [("unique ID", False)], 9.5, GREY, align=PP_ALIGN.CENTER)

chip(7.35, DTOP + 0.34, 4.98, 0.62,
     "1 booking  ·  one round trip, one purchase decision",
     NAVY, WHITE, 13.0, bold=True)
textbox(7.35, DTOP + 1.06, 4.98, 0.5,
        [("One row, one trip purpose — the grain every rule, chart and "
          "segment in this deck is built on.", False)], 10.5, GREY, spacing=1.2)

# ── closer, in the pattern the other section slides use ───────────────────────
bar(1.0, 6.35, 0.04, 0.26, BRIGHT)
textbox(1.16, 6.33, 10.84, 0.4,
        [("38,116,259 coupons became ", False), ("22,911,450 bookings", True),
         (" — one purchase decision, one trip purpose — rolling up again to 13.4M customers.",
          False)], 12.5, GREY)

# ── move it to the head of the EDA section (becomes slide 5) ──────────────────
sld_ids = prs.slides._sldIdLst
new = sld_ids[-1]
sld_ids.remove(new)
sld_ids.insert(4, new)

prs.save(DST)
print(f"saved {DST}  ({len(prs.slides)} slides)")
